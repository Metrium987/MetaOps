"""Multi-format file parsers for RAG ingestion.

Supports: PDF, DOCX, DOC, PPTX, PPT, RTF, HTML, JSON, JSONL, TXT, Markdown, XLS, XLSX.
Each parser extracts text content from its respective format.
"""

import json
import logging
import re
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)


def parse_file(file_path: str, content: Optional[bytes] = None) -> str:
    """Parse a file and extract text content based on its extension.

    Args:
        file_path: Path to the file.
        content: Raw file bytes (optional, will read from file_path if not provided).

    Returns:
        Extracted text content, or empty string if parsing fails.
    """
    path = Path(file_path)
    suffix = path.suffix.lower()

    parser_map = {
        ".pdf": _parse_pdf,
        ".docx": _parse_docx,
        ".doc": _parse_doc,
        ".pptx": _parse_pptx,
        ".ppt": _parse_ppt,
        ".rtf": _parse_rtf,
        ".html": _parse_html,
        ".htm": _parse_html,
        ".json": _parse_json,
        ".jsonl": _parse_jsonl,
        ".txt": _parse_text,
        ".md": _parse_markdown,
        ".csv": _parse_text,
        ".log": _parse_text,
        ".xml": _parse_xml,
        ".yaml": _parse_yaml,
        ".yml": _parse_yaml,
        ".toml": _parse_toml,
        ".ini": _parse_ini,
        ".cfg": _parse_ini,
        ".py": _parse_code,
        ".js": _parse_code,
        ".ts": _parse_code,
        ".java": _parse_code,
        ".go": _parse_code,
        ".rs": _parse_code,
        ".c": _parse_code,
        ".cpp": _parse_code,
        ".h": _parse_code,
        ".hpp": _parse_code,
        ".rb": _parse_code,
        ".php": _parse_code,
        ".sh": _parse_code,
        ".bash": _parse_code,
        ".sql": _parse_code,
        ".r": _parse_code,
        ".scala": _parse_code,
        ".kt": _parse_code,
        ".swift": _parse_code,
    }

    parser = parser_map.get(suffix)
    if parser is None:
        logger.warning("No parser for file type '%s', falling back to text extraction", suffix)
        parser = _parse_text

    try:
        if content is None:
            content = path.read_bytes()
        return parser(content)
    except Exception as e:
        logger.warning("Failed to parse %s: %s", file_path, e)
        return ""


# ── PDF Parser ────────────────────────────────────────────────────────────────

def _parse_pdf(content: bytes) -> str:
    """Extract text from PDF using PyMuPDF (fitz) or fallback to pdfplumber."""
    try:
        import fitz
        doc = fitz.open(stream=content, filetype="pdf")
        text_parts = []
        for page in doc:
            text_parts.append(page.get_text())
        doc.close()
        return "\n".join(text_parts)
    except ImportError:
        pass

    try:
        import pdfplumber
        import io
        with pdfplumber.open(io.BytesIO(content)) as pdf:
            text_parts = []
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
            return "\n".join(text_parts)
    except ImportError:
        pass

    logger.warning("No PDF parser available (install PyMuPDF or pdfplumber)")
    return ""


# ── DOCX Parser ───────────────────────────────────────────────────────────────

def _parse_docx(content: bytes) -> str:
    """Extract text from DOCX using python-docx."""
    try:
        from docx import Document
        import io
        doc = Document(io.BytesIO(content))
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    text_parts.append(row_text)
        return "\n".join(text_parts)
    except ImportError:
        logger.warning("No DOCX parser available (install python-docx)")
        return ""


# ── DOC Parser ────────────────────────────────────────────────────────────────

def _parse_doc(content: bytes) -> str:
    """Extract text from old .doc format using antiword or catdoc."""
    import subprocess
    import tempfile
    import os

    for cmd in [["antiword", "-"], ["catdoc", "-"]]:
        try:
            with tempfile.NamedTemporaryFile(suffix=".doc", delete=False) as tmp:
                tmp.write(content)
                tmp_path = tmp.name
            try:
                result = subprocess.run(
                    [cmd[0], tmp_path] if cmd[0] == "antiword" else [cmd[0], tmp_path],
                    capture_output=True, text=True, timeout=30
                )
                if result.returncode == 0 and result.stdout.strip():
                    return result.stdout
            finally:
                os.unlink(tmp_path)
        except (FileNotFoundError, subprocess.TimeoutExpired):
            continue

    logger.warning("No DOC parser available (install antiword or catdoc)")
    return ""


# ── PPTX Parser ──────────────────────────────────────────────────────────────

def _parse_pptx(content: bytes) -> str:
    """Extract text from PPTX using python-pptx."""
    try:
        from pptx import Presentation
        import io
        prs = Presentation(io.BytesIO(content))
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
        return "\n".join(text_parts)
    except ImportError:
        logger.warning("No PPTX parser available (install python-pptx)")
        return ""


# ── PPT Parser ────────────────────────────────────────────────────────────────

def _parse_ppt(content: bytes) -> str:
    """Extract text from old .ppt using catppt or LibreOffice conversion."""
    import subprocess
    import tempfile
    import os

    try:
        with tempfile.NamedTemporaryFile(suffix=".ppt", delete=False) as tmp:
            tmp.write(content)
            tmp_path = tmp.name
        try:
            result = subprocess.run(
                ["catppt", tmp_path], capture_output=True, text=True, timeout=30
            )
            if result.returncode == 0 and result.stdout.strip():
                return result.stdout
        finally:
            os.unlink(tmp_path)
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass

    logger.warning("No PPT parser available (install catppt)")
    return ""


# ── RTF Parser ────────────────────────────────────────────────────────────────

def _parse_rtf(content: bytes) -> str:
    """Extract text from RTF using striprtf or unrtf."""
    try:
        from striprtf.striprtf import rtf_to_text
        return rtf_to_text(content.decode("utf-8", errors="replace"))
    except ImportError:
        pass

    import subprocess
    import tempfile
    import os

    try:
        with tempfile.NamedTemporaryFile(suffix=".rtf", delete=False) as tmp:
            tmp.write(content)
            tmp_path = tmp.name
        try:
            result = subprocess.run(
                ["unrtf", "--text", tmp_path], capture_output=True, text=True, timeout=30
            )
            if result.returncode == 0 and result.stdout.strip():
                return result.stdout
        finally:
            os.unlink(tmp_path)
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass

    logger.warning("No RTF parser available (install striprtf or unrtf)")
    return ""


# ── HTML Parser ───────────────────────────────────────────────────────────────

def _parse_html(content: bytes) -> str:
    """Extract text from HTML using BeautifulSoup or html2text."""
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(content, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        return soup.get_text(separator="\n", strip=True)
    except ImportError:
        pass

    try:
        import html2text
        h = html2text.HTML2Text()
        h.ignore_links = False
        h.ignore_images = True
        return h.handle(content.decode("utf-8", errors="replace"))
    except ImportError:
        pass

    text = content.decode("utf-8", errors="replace")
    text = re.sub(r"<script[^>]*>.*?</script>", "", text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r"<style[^>]*>.*?</style>", "", text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r"<[^>]+>", " ", text)
    text = re.sub(r"\s+", " ", text).strip()
    return text


# ── JSON Parser ───────────────────────────────────────────────────────────────

def _parse_json(content: bytes) -> str:
    """Extract text from JSON, flattening nested structures."""
    try:
        data = json.loads(content.decode("utf-8", errors="replace"))
        return _flatten_json(data)
    except json.JSONDecodeError:
        return content.decode("utf-8", errors="replace")


def _flatten_json(obj, prefix="") -> str:
    """Recursively flatten JSON into readable text."""
    parts = []
    if isinstance(obj, dict):
        for k, v in obj.items():
            parts.append(_flatten_json(v, f"{prefix}{k}: " if prefix else f"{k}: "))
    elif isinstance(obj, list):
        for i, item in enumerate(obj):
            parts.append(_flatten_json(item, f"{prefix}[{i}] "))
    else:
        parts.append(f"{prefix}{obj}")
    return "\n".join(p for p in parts if p.strip())


# ── JSONL Parser ──────────────────────────────────────────────────────────────

def _parse_jsonl(content: bytes) -> str:
    """Extract text from JSONL (one JSON object per line)."""
    text = content.decode("utf-8", errors="replace")
    parts = []
    for line in text.strip().split("\n"):
        line = line.strip()
        if not line:
            continue
        try:
            obj = json.loads(line)
            parts.append(_flatten_json(obj))
        except json.JSONDecodeError:
            parts.append(line)
    return "\n".join(parts)


# ── Text/Markdown/Code Parsers ────────────────────────────────────────────────

def _parse_text(content: bytes) -> str:
    return content.decode("utf-8", errors="replace")


def _parse_markdown(content: bytes) -> str:
    text = content.decode("utf-8", errors="replace")
    text = re.sub(r"```[\s\S]*?```", "[code block]", text)
    text = re.sub(r"`[^`]+`", "", text)
    text = re.sub(r"!?\[([^\]]*)\]\([^)]*\)", r"\1", text)
    text = re.sub(r"^#{1,6}\s+", "", text, flags=re.MULTILINE)
    text = re.sub(r"\*{1,2}([^*]+)\*{1,2}", r"\1", text)
    return text.strip()


def _parse_code(content: bytes) -> str:
    return content.decode("utf-8", errors="replace")


# ── XML Parser ────────────────────────────────────────────────────────────────

def _parse_xml(content: bytes) -> str:
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(content, "xml")
        return soup.get_text(separator="\n", strip=True)
    except ImportError:
        pass

    text = content.decode("utf-8", errors="replace")
    text = re.sub(r"<[^>]+>", " ", text)
    text = re.sub(r"\s+", " ", text).strip()
    return text


# ── YAML Parser ───────────────────────────────────────────────────────────────

def _parse_yaml(content: bytes) -> str:
    try:
        import yaml
        data = yaml.safe_load(content.decode("utf-8", errors="replace"))
        return _flatten_json(data) if data else ""
    except ImportError:
        return content.decode("utf-8", errors="replace")


# ── TOML Parser ───────────────────────────────────────────────────────────────

def _parse_toml(content: bytes) -> str:
    try:
        import tomllib
        data = tomllib.loads(content.decode("utf-8", errors="replace"))
        return _flatten_json(data)
    except ImportError:
        pass

    try:
        import tomli
        data = tomli.loads(content.decode("utf-8", errors="replace"))
        return _flatten_json(data)
    except ImportError:
        return content.decode("utf-8", errors="replace")


# ── INI Parser ────────────────────────────────────────────────────────────────

def _parse_ini(content: bytes) -> str:
    import configparser
    import io
    config = configparser.ConfigParser()
    config.read_string(content.decode("utf-8", errors="replace"))
    parts = []
    for section in config.sections():
        parts.append(f"[{section}]")
        for key, value in config.items(section):
            parts.append(f"{key} = {value}")
    return "\n".join(parts)
